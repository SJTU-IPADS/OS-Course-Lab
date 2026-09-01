"""LaTeX math becomes a native PowerPoint equation (OMML) in the pptx export."""

import tempfile
import unittest
from pathlib import Path

from lxml import etree
from pptx import Presentation

from lecturekit import Lecture
from lecturekit.renderers.pptx import PptxRenderer
from lecturekit.renderers.pptx.omml import M, plain_text, to_omml, tokenize
from lecturekit.renderers.pptx.text import parse_markdown


def xml(latex, **kwargs):
    return etree.tostring(to_omml(latex, **kwargs), encoding="unicode")


def tags(latex):
    """Local names of the OMML nodes, so a test can assert on structure."""
    return [etree.QName(n).localname for n in to_omml(latex).iter()]


class TokenizeTest(unittest.TestCase):
    def test_splits_commands_and_characters(self):
        self.assertEqual(tokenize(r"x^2"), ["x", "^", "2"])
        self.assertEqual(tokenize(r"\frac{a}{b}"),
                         ["\\frac", "{", "a", "}", "{", "b", "}"])

    def test_eats_the_space_that_ends_a_control_word(self):
        # `\times D` is TeX for two symbols with no gap between them.
        self.assertEqual(tokenize(r"B\times D"), ["B", "\\times", "D"])

    def test_keeps_a_space_between_plain_characters(self):
        self.assertEqual(tokenize("a = b"), ["a", " ", "=", " ", "b"])


class StructureTest(unittest.TestCase):
    def test_subscript_and_superscript(self):
        self.assertIn("sSub", tags(r"x_i"))
        self.assertIn("sSup", tags(r"x^2"))

    def test_both_scripts_fold_into_one_base(self):
        structure = tags(r"C_{a}^{b}")
        self.assertIn("sSubSup", structure)
        self.assertNotIn("sSub", structure)

    def test_fraction_and_root(self):
        self.assertIn("f", tags(r"\frac{a}{b}"))
        self.assertIn("rad", tags(r"\sqrt{x}"))

    def test_big_operator_carries_its_limits(self):
        out = xml(r"\sum_{i=1}^{N}x_i")
        self.assertIn("nary", tags(r"\sum_{i=1}^{N}x_i"))
        self.assertIn("∑", out)
        self.assertIn("undOvr", out)  # limits stack above and below

    def test_the_sum_body_stops_at_a_relation(self):
        # `O(i)` belongs under the ∑; `= O(N)` does not.
        node = to_omml(r"\sum_{i=1}^{N}O(i) = O(N)")
        nary = node.find(f"{{{M}}}nary")
        body = "".join(t.text for t in nary.find(f"{{{M}}}e").iter(f"{{{M}}}t"))
        self.assertEqual(body, "O(i)")

    def test_left_right_becomes_a_sized_delimiter(self):
        out = xml(r"\left(\frac{a}{b}\right)")
        self.assertIn("d", tags(r"\left(\frac{a}{b}\right)"))
        self.assertIn('begChr', out)

    def test_a_group_inside_a_delimiter_does_not_swallow_the_closer(self):
        # Regression: the sum must not eat \right, leaving an unclosed paren.
        out = xml(r"\left(\sum_{i} x_i\right) + 1")
        self.assertIn("endChr", out)
        self.assertTrue(out.rstrip().endswith("</m:oMath>"))

    def test_function_names_are_upright(self):
        out = xml(r"\operatorname{softmax}(x)")
        self.assertIn('<m:sty m:val="p"/>', out)
        self.assertIn("softmax", out)

    def test_blackboard_bold_is_double_struck_and_upright(self):
        out = xml(r"\mathbb{R}")
        self.assertIn('double-struck', out)
        self.assertIn('<m:sty m:val="p"/>', out)

    def test_symbols_become_their_unicode_character(self):
        self.assertIn("×", xml(r"B\times D"))
        self.assertIn("≈", xml(r"a\approx b"))

    def test_an_unknown_command_keeps_its_name_rather_than_vanishing(self):
        self.assertIn("widehat", xml(r"\widehat{x}"))

    def test_display_math_is_centered(self):
        self.assertIn('<m:jc m:val="center"/>', xml(r"x=1", display=True))

    def test_plain_text_reads_as_the_formula(self):
        self.assertEqual(plain_text(r"O(N^2)"), "O(N^2)")
        self.assertEqual(plain_text(r"B\times D"), "B×D")


class ParseTest(unittest.TestCase):
    def test_inline_math_becomes_a_math_run(self):
        [para] = parse_markdown("每次 forward 是 $O(N)$")
        self.assertEqual([r.math for r in para.runs][-1], "O(N)")

    def test_math_inside_an_autobolded_headline_is_still_parsed(self):
        # The DSL bolds a flush-left headline, wrapping any $…$ with it.
        [para] = parse_markdown("**第 $i$ 次 forward**")
        self.assertIn("i", [r.math for r in para.runs])

    def test_display_math_is_its_own_paragraph(self):
        # `$$` is indented by one space: the DSL's escape from auto-bold, which
        # would otherwise wrap the fence itself in asterisks.
        paras = parse_markdown("前文\n $$\n C = O(N)\n $$\n后文")
        self.assertEqual([p.kind for p in paras], ["para", "math", "para"])
        self.assertEqual(paras[1].runs[0].math, "C = O(N)")

    def test_a_multi_line_display_formula_stays_one_paragraph(self):
        paras = parse_markdown(" $$\n C\n = O(N)\n $$")
        self.assertEqual(len(paras), 1)
        self.assertEqual(paras[0].runs[0].math, "C = O(N)")


class RenderTest(unittest.TestCase):
    def deck(self, body_md):
        lecture = Lecture(id="lec", title="Deck")

        def body(p):
            p.title("Math")
            p.slide(body_md)

        lecture.page(id="math", body=body)
        tmp = Path(tempfile.mkdtemp())
        out = PptxRenderer().render(lecture.build(), tmp)
        return Presentation(str(out))

    def slide_xml(self, deck):
        return "\n".join(s._element.xml for s in deck.slides[0].shapes)

    def test_equation_is_embedded_not_printed_as_latex(self):
        deck = self.deck("成本是 $O(N^2)$")
        text = "\n".join(s.text_frame.text for s in deck.slides[0].shapes
                         if s.has_text_frame)
        self.assertNotIn("$", text)
        self.assertIn("oMath", self.slide_xml(deck))

    def test_equation_has_a_plain_text_fallback(self):
        deck = self.deck(" $$\n O(N^2)\n $$")
        self.assertIn("Fallback", self.slide_xml(deck))
        self.assertIn("O(N^2)", self.slide_xml(deck))

    def test_math_runs_carry_the_body_size(self):
        from lecturekit.renderers.pptx import theme
        deck = self.deck("成本是 $O(N)$")
        self.assertIn(f'sz="{theme.BODY_PT * 100}"', self.slide_xml(deck))


if __name__ == "__main__":
    unittest.main()
