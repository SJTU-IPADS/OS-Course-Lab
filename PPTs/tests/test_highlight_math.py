"""Math inside a `p.highlight(...)` chip, in each of the four targets.

A punchline is often the formula the page spent ten slides earning, so a chip
takes `$…$` and `$$…$$` like slide text does. What each renderer has to say is
where display math *goes* when the chip is one line: nowhere, in every case —
it stays on its row, set in display style.
"""

import tempfile
import unittest
from pathlib import Path

from lxml import etree
from pptx import Presentation

from lecturekit import model
from lecturekit.dsl import Lecture
from lecturekit.renderers.latex.assets import AssetCopier
from lecturekit.renderers.latex.blocks import Ctx, emit_block
from lecturekit.renderers.pptx import PptxRenderer
from lecturekit.renderers.transcript import build_html
from lecturekit.renderers.transcript.images import Embedder
from lecturekit.renderers.viewer import build_marp_markdown

FENCED = "代价\n$$\n\\frac{N}{W}\n$$"


def deck(text, *, tone="yellow", footnotes=()):
    lecture = Lecture(id="lec01", title="L")

    def body(p):
        p.title("T")
        chip = p.highlight(text, tone=tone)
        for note in footnotes:
            chip.footnote(note)

    lecture.page("p1", body=body)
    return build_marp_markdown(lecture.build())


class HighlightLinesTest(unittest.TestCase):
    def test_a_fence_folds_onto_one_row(self):
        self.assertEqual(
            model.highlight_lines(FENCED), ["代价", "$$\\frac{N}{W}$$"]
        )

    def test_a_multi_line_formula_joins_with_spaces(self):
        self.assertEqual(
            model.highlight_lines("$$\na = b\n+ c\n$$"), ["$$a = b + c$$"]
        )

    def test_a_one_line_formula_is_already_a_row(self):
        self.assertEqual(model.highlight_lines("$$x^2$$"), ["$$x^2$$"])

    def test_rows_are_stripped_like_any_chip(self):
        self.assertEqual(
            model.highlight_lines("  Underloaded  \n 一半的队列是空的 "),
            ["Underloaded", "一半的队列是空的"],
        )


class ViewerTest(unittest.TestCase):
    def test_a_chip_without_math_is_still_a_raw_html_paragraph(self):
        self.assertIn(
            '<p class="lk-highlight lk-highlight--yellow">'
            "<span>Underloaded</span></p>",
            deck("Underloaded"),
        )

    def test_math_switches_the_chip_to_a_markdown_line(self):
        # Marp never looks inside a raw HTML block, so a chip carrying math is
        # drawn on a <span>: text follows it on the line, so it is a paragraph.
        self.assertIn(
            '<span class="lk-highlight lk-highlight--orange">'
            "<span>每步代价 $T \\approx 2PN/W$</span></span>",
            deck("每步代价 $T \\approx 2PN/W$", tone="orange"),
        )

    def test_a_display_row_is_inline_math_in_display_style(self):
        self.assertIn(
            "<span>代价<br>$\\displaystyle \\frac{N}{W}$</span>", deck(FENCED)
        )

    def test_the_formula_reaches_marp_unescaped(self):
        # `&lt;` is not what MathJax should typeset.
        self.assertIn("$a_i < b^2$", deck("$a_i < b^2$"))

    def test_markdown_around_the_formula_cannot_fire_twice(self):
        # The chip's own inline markdown has already run; what is left of it
        # must not be read a second time when Marp parses the line.
        md = deck("$x$ 和 **粗** 和 snake_case 和 a * b")
        self.assertIn("<strong>粗</strong>", md)
        self.assertIn("snake&#95;case", md)
        self.assertIn("a &#42; b", md)

    def test_a_footnote_marker_still_tucks_inside_the_chip(self):
        md = deck(FENCED, footnotes=["λ < μ"])
        self.assertIn('</span> <sup class="footnote-ref">1</sup></span>', md)
        self.assertNotIn("</span></span> <sup", md)


class LatexTest(unittest.TestCase):
    def emit(self, text, **kw):
        tmp = Path(tempfile.mkdtemp())
        ctx = Ctx(
            lecture_id="lec-a",
            page_id="p1",
            slide_width=1280,
            assets=AssetCopier(tmp / "out"),
            asset_root=tmp / "src",
        )
        block = model.Block(
            kind="highlight", content={"text": text, "tone": "yellow"}, **kw
        )
        return emit_block(block, ctx)

    def test_a_display_row_is_set_in_display_style(self):
        # A chip's rows are the cells of a tabular, and TeX cannot open display
        # math inside one.
        self.assertEqual(
            self.emit(FENCED),
            r"\lkhighlight{lkChipYellow}{代价 \\ $\displaystyle \frac{N}{W}$}",
        )

    def test_inline_math_passes_through_as_tex(self):
        self.assertIn(
            r"$T \approx 2PN/W$", self.emit(r"每步代价 $T \approx 2PN/W$")
        )


class TranscriptTest(unittest.TestCase):
    def html(self, text):
        lecture = Lecture(id="lec-t", title="T")
        lecture.page(
            "p1", body=lambda p: (p.title("T"), p.highlight(text))
        )
        built = lecture.build()
        return build_html(built, Embedder(None, built.borrowed))

    def test_a_display_row_is_typeset_in_the_sheet_s_own_math(self):
        html = self.html(FENCED)
        self.assertIn('<p class="tx-highlight tx-yellow">', html)
        self.assertIn('<span class="tx-frac">', html)
        self.assertNotIn("$$", html)

    def test_inline_math_still_works(self):
        self.assertNotIn("$", self.html(r"每步代价 $T \approx 2PN/W$"))


class PptxTest(unittest.TestCase):
    def chip_xml(self, text):
        lecture = Lecture(id="lec", title="Deck")
        lecture.page("p1", body=lambda p: (p.title("T"), p.highlight(text)))
        tmp = Path(tempfile.mkdtemp())
        out = PptxRenderer().render(lecture.build(), tmp)
        slide = Presentation(str(out)).slides[0]
        shape = next(
            s for s in slide.shapes
            if s.has_text_frame and "代价" in s.text_frame.text
        )
        return etree.tostring(shape.text_frame._txBody, encoding="unicode")

    def test_a_display_row_becomes_a_native_equation(self):
        xml = self.chip_xml(FENCED)
        # PowerPoint's own equation markup — a fraction, not the LaTeX source.
        self.assertIn("oMath", xml)
        self.assertIn("}m:f", xml.replace("<", "}"))
        self.assertNotIn("$$", xml)


if __name__ == "__main__":
    unittest.main()
