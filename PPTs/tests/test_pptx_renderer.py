import base64
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Inches, Pt

from lecturekit import Lecture
from lecturekit.renderers import RENDERERS, get_renderer
from lecturekit.renderers.pptx import PptxRenderer, theme
from lecturekit.renderers.pptx.layout import Layout

# 1x1 transparent PNG (python-pptx can embed it; svg is unsupported upstream).
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9Q"
    "DwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


def shape_texts(slide):
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame]


def all_text(slide):
    return "\n".join(shape_texts(slide))


class PptxRendererTest(unittest.TestCase):
    def render(self, lecture, *, asset_root=None):
        tmp = Path(tempfile.mkdtemp())
        out = PptxRenderer(asset_root=asset_root).render(lecture.build(), tmp)
        return Presentation(str(out)), out

    def sample(self):
        lecture = Lecture(id="lec", title="Deck Title")

        def body(p):
            p.title("Welcome")
            p.slide("- hello\n- **world**")
            p.code("python", "print('hi')")
            p.link("OSTEP", "https://example.test")
            p.table([["a", "b"]], headers=["H1", "H2"])
            p.aside("an aside")
            p.sidenote("Ref", "see this")
            p.notes("teacher only")

        with lecture.section("Intro", id="intro") as s:
            s.page("welcome", body=body)
        return lecture

    def test_registered_as_render_target(self):
        self.assertIn("pptx", RENDERERS)
        self.assertIs(get_renderer("pptx"), PptxRenderer)

    def test_one_slide_per_page(self):
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.slide("x")

        with lecture.section("S", id="s") as s:
            s.page("p1", body=body)
            s.page("p2", body=body)
        prs, _ = self.render(lecture)
        self.assertEqual(len(prs.slides), 2)

    def test_output_named_after_lecture_title(self):
        _, out = self.render(self.sample())
        self.assertEqual(out.name, "deck-title.pptx")

    def test_slide_size_follows_ratio(self):
        lecture = Lecture(id="lec", title="L", ratio="4:3")

        def body(p):
            p.title("P")
            p.slide("x")

        with lecture.section("S", id="s") as s:
            s.page("p1", body=body)
        prs, _ = self.render(lecture)
        self.assertEqual(prs.slide_width, Emu(Inches(960 / 96)))

    def test_title_and_bullet_text_present(self):
        prs, _ = self.render(self.sample())
        text = all_text(prs.slides[0])
        self.assertIn("Welcome", text)
        self.assertIn("hello", text)
        self.assertIn("world", text)

    def test_cover_slide_draws_title_metadata_and_logos(self):
        src = Path(tempfile.mkdtemp())
        (src / "assets").mkdir()
        (src / "assets" / "left.png").write_bytes(_PNG)
        (src / "assets" / "right.png").write_bytes(_PNG)
        lecture = Lecture(id="lec", title="L")
        lecture.cover(
            "Elastic model serving via efficient autoscaling",
            author="Xingda Wei",
            time="July 2026",
            logo=("assets/left.png", "assets/right.png"),
        )

        prs, _ = self.render(lecture, asset_root=src)
        slide = prs.slides[0]
        text = all_text(slide)

        self.assertIn("Elastic model serving via efficient autoscaling", text)
        self.assertIn("Xingda Wei", text)
        self.assertIn("July 2026", text)
        pics = [
            sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertEqual(len(pics), 2)

    def test_runs_carry_east_asian_typeface(self):
        # font.name only sets the Latin face; CJK text needs an <a:ea> typeface
        # or PowerPoint substitutes its default CJK font (format mismatch).
        from pptx.oxml.ns import qn
        prs, _ = self.render(self.sample())
        eas = []
        for shape in prs.slides[0].shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    ea = run.font._rPr.find(qn("a:ea"))
                    if ea is not None:
                        eas.append(ea.get("typeface"))
        self.assertTrue(eas, "no run set an East Asian typeface")
        self.assertIn("PingFang SC", eas)

    def test_title_rule_is_thin_but_visible(self):
        # visible (not a hairline) yet not a heavy banner
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Pt
        prs, _ = self.render(self.sample())
        rules = [
            s for s in prs.slides[0].shapes
            if s.shape_type == MSO_SHAPE_TYPE.LINE
        ]
        self.assertEqual(len(rules), 1)
        self.assertGreaterEqual(rules[0].line.width, Pt(1))
        self.assertLessEqual(rules[0].line.width, Pt(2))

    def test_slide_body_uses_exact_line_spacing_matching_the_estimate(self):
        # exact point spacing (size * 1.4), not the 1.4 multiple — PowerPoint
        # multiplies a float by the font's leading, rendering taller than the
        # estimate and overlapping the next block
        from pptx.util import Pt
        prs, _ = self.render(self.sample())
        body = max(
            (s for s in prs.slides[0].shapes if s.has_text_frame),
            key=lambda s: len(s.text_frame.paragraphs),
        )
        for p in body.text_frame.paragraphs:
            self.assertEqual(p.line_spacing, Pt(20 * 1.4))   # sample body is 20pt

    def test_text_boxes_have_zero_vertical_margin(self):
        # auto-fit grows a box by its margins; nonzero margins make the grown
        # box overlap the block below
        prs, _ = self.render(self.sample())
        body = max(
            (s for s in prs.slides[0].shapes if s.has_text_frame),
            key=lambda s: len(s.text_frame.paragraphs),
        )
        self.assertEqual(body.text_frame.margin_top, 0)
        self.assertEqual(body.text_frame.margin_bottom, 0)

    def test_notes_block_is_excluded(self):
        prs, _ = self.render(self.sample())
        self.assertNotIn("teacher only", all_text(prs.slides[0]))

    def test_code_text_present(self):
        prs, _ = self.render(self.sample())
        self.assertIn("print('hi')", all_text(prs.slides[0]))

    def test_a_demo_reaches_the_pptx_as_a_still_transcript(self):
        # A PPTX is a file on someone else's laptop: nothing in it runs, so the
        # block has to print what the deck prints when nobody presses -- and it
        # has to print it at all, which it did not before demos were runnable.
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.demo("see the assembly", "gcc -S a.c", output="a.s written")

        lecture.page("p1", body=body)
        prs, _ = self.render(lecture)
        text = all_text(prs.slides[0])
        self.assertIn("see the assembly", text)
        self.assertIn("$ gcc -S a.c", text)
        self.assertIn("a.s written", text)

    def test_code_background_matches_rows_without_an_extra_blank_line(self):
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.code("text", "one\ntwo\nthree")

        lecture.page("p1", body=body)
        prs, _ = self.render(lecture)
        code_box = next(
            s for s in prs.slides[0].shapes
            if s.has_text_frame and s.text_frame.text == "one\ntwo\nthree"
        )
        expected = Pt(3 * theme.CODE_PT * 1.4 + 8)
        self.assertEqual(code_box.height, expected)
        self.assertEqual(code_box.text_frame.margin_top, Pt(4))
        self.assertEqual(code_box.text_frame.margin_bottom, Pt(4))
        for paragraph in code_box.text_frame.paragraphs:
            self.assertAlmostEqual(
                paragraph.line_spacing, Pt(theme.CODE_PT * 1.4), delta=Pt(0.02)
            )

    def test_page_gap_spreads_pptx_block_groups(self):
        def make(fill):
            lecture = Lecture(id="lec", title="L")

            def body(p):
                p.title("P")
                if fill:
                    p.gap("fill")
                p.slide("first")
                p.slide("second")

            lecture.page("p1", body=body)
            return lecture

        normal, _ = self.render(make(False))
        spread, _ = self.render(make(True))

        def top_of(prs, text):
            return next(
                s.top for s in prs.slides[0].shapes
                if s.has_text_frame and s.text_frame.text == text
            )

        self.assertGreater(top_of(spread, "second"), top_of(normal, "second"))
        self.assertEqual(top_of(spread, "first"), top_of(normal, "first"))

    def test_link_carries_hyperlink(self):
        prs, _ = self.render(self.sample())
        urls = [
            run.hyperlink.address
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        self.assertIn("https://example.test", urls)

    def test_table_rendered_as_native_table(self):
        prs, _ = self.render(self.sample())
        tables = [
            s for s in prs.slides[0].shapes if s.has_table
        ]
        self.assertEqual(len(tables), 1)
        table = tables[0].table
        self.assertEqual(len(table.rows), 2)   # header + 1 data row
        self.assertEqual(len(table.columns), 2)
        self.assertEqual(table.cell(0, 0).text, "H1")
        self.assertEqual(table.cell(1, 1).text, "b")

    def test_sidenote_and_aside_text_present(self):
        prs, _ = self.render(self.sample())
        text = all_text(prs.slides[0])
        self.assertIn("Ref", text)
        self.assertIn("see this", text)
        self.assertIn("an aside", text)

    def test_sidenote_image_logo_is_embedded_not_shown_as_path(self):
        src = Path(tempfile.mkdtemp())
        (src / "assets").mkdir()
        (src / "assets" / "logo.png").write_bytes(_PNG)
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.sidenote("Book", "a note", logo="assets/logo.png")

        with lecture.section("S", id="s") as s:
            s.page("p1", body=body)
        prs, _ = self.render(lecture, asset_root=src)
        slide = prs.slides[0]
        # the path must not leak into the text, and the logo must be a picture
        self.assertNotIn("assets/logo.png", all_text(slide))
        pics = [
            sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertEqual(len(pics), 1)

    def test_image_embedded_as_picture(self):
        src = Path(tempfile.mkdtemp())
        (src / "assets").mkdir()
        (src / "assets" / "pic.png").write_bytes(_PNG)
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.image("assets/pic.png", caption="a cap")

        with lecture.section("S", id="s") as s:
            s.page("p1", body=body)
        prs, _ = self.render(lecture, asset_root=src)
        pics = [
            s for s in prs.slides[0].shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertEqual(len(pics), 1)
        self.assertIn("a cap", all_text(prs.slides[0]))

    def test_explicit_image_height_is_honored_without_distortion(self):
        src = Path(tempfile.mkdtemp())
        (src / "assets").mkdir()
        (src / "assets" / "pic.png").write_bytes(_PNG)
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.image("assets/pic.png", height_px=180)

        lecture.page("p1", body=body)
        prs, _ = self.render(lecture, asset_root=src)
        pic = next(
            s for s in prs.slides[0].shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        )
        self.assertEqual(pic.height, Inches(180 / 96))
        self.assertEqual(pic.width, pic.height)  # source fixture is square

    def test_highlight_is_a_centered_toned_chip(self):
        from pptx.enum.text import PP_ALIGN

        lecture = Lecture(id="lec", title="Deck")
        lecture.page(
            "p1", body=lambda p: (p.title("T"), p.highlight("Underloaded")),
        )
        prs, _ = self.render(lecture)
        slide = prs.slides[0]
        chips = [s for s in slide.shapes if s.has_text_frame
                 and s.text_frame.text == "Underloaded"]
        self.assertEqual(len(chips), 1)
        chip = chips[0]
        para = chip.text_frame.paragraphs[0]
        self.assertEqual(para.alignment, PP_ALIGN.CENTER)
        self.assertTrue(para.runs[0].font.bold)
        self.assertEqual(chip.fill.type, MSO_FILL.BACKGROUND)
        self.assertEqual(para.runs[0].font.color.rgb, theme.CHIP_INK["yellow"])
        # the frame carries the tone too — it is what makes the chip read as one
        # object with nothing painted inside it
        self.assertEqual(chip.line.color.rgb, theme.CHIP_INK["yellow"])

    def test_highlight_chip_is_narrower_than_the_content_column(self):
        lecture = Lecture(id="lec", title="Deck")
        lecture.page("p1", body=lambda p: (p.title("T"), p.highlight("短")))
        prs, _ = self.render(lecture)
        chip = next(s for s in prs.slides[0].shapes if s.has_text_frame
                    and s.text_frame.text == "短")
        layout = Layout.from_ratio("16:9")
        self.assertLess(chip.width, layout.content_width)
        self.assertGreater(chip.width, 0)

    def test_highlight_uses_one_paragraph_per_line(self):
        lecture = Lecture(id="lec", title="Deck")
        lecture.page(
            "p1",
            body=lambda p: (p.title("T"), p.highlight("a\nb", tone="orange")),
        )
        prs, _ = self.render(lecture)
        chip = next(s for s in prs.slides[0].shapes if s.has_text_frame
                    and s.text_frame.text.replace("\n", "") == "ab")
        self.assertEqual(len(chip.text_frame.paragraphs), 2)
        self.assertEqual(
            chip.text_frame.paragraphs[0].runs[0].font.color.rgb,
            theme.CHIP_INK["orange"],
        )

    def test_deferred_blocks_are_skipped_without_error(self):
        lecture = Lecture(id="lec", title="L")

        def body(p):
            p.title("P")
            p.slide("text")
            arch = p.architecture()
            arch.layer("L1", ["m1", "m2"])

        with lecture.section("S", id="s") as s:
            s.page("p1", body=body)
        prs, _ = self.render(lecture)   # must not raise
        self.assertEqual(len(prs.slides), 1)


if __name__ == "__main__":
    unittest.main()
