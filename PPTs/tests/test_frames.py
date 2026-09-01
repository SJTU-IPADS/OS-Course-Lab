"""``p.frames(...)``: one page, one figure per frame — expansion and its fallout."""

import contextlib
import io
import tempfile
import unittest
from pathlib import Path

from lecturekit import model, references
from lecturekit.book import BookModel
from lecturekit.cli import format_tree
from lecturekit.dsl import Lecture
from lecturekit.model import ValidationError
from lecturekit.renderers.latex import LatexRenderer, coverage
from lecturekit.renderers.viewer import (
    build_data,
    build_marp_markdown,
    build_outline_html,
)
from lecturekit.renderers.viewer.pdf import SLIDE_LINK_PREFIX
from lecturekit.serialize import lecture_to_dict

SOURCES = ("log-1.png", "log-2.png", "log-3.png")
NOTE = "讲：三帧对应写 log、commit、改数据"


def _animation(p):
    p.title("Commit logging 先写完整新状态")
    p.slide("log 先落盘，再改 bank file")
    p.frames(
        *SOURCES,
        width_px=900,
        caption="先写 log，再改数据",
    ).footnote("Saltzer & Kaashoek §9.3")
    p.notes(NOTE)


def _lecture(**page_kwargs):
    """before / a three-frame animation / after — deck order 1..5."""
    lec = Lecture(id="lec-f", title="Lecture F")

    def before(p):
        p.title("Before")
        p.slide("first")

    def after(p):
        p.title("After")
        p.slide("last")

    with lec.section("S") as s:
        s.page(id="before", body=before)
        s.page(id="anim", body=_animation, **page_kwargs)
        s.page(id="after", body=after)
    return lec.build()


def _pages(lecture):
    return model.flatten_pages(lecture.children)


def _page_ids(lecture):
    return [page.id for page in _pages(lecture)]


def _images(page):
    return [block for block in page.blocks if block.kind == "image"]


class ExpansionTest(unittest.TestCase):
    def test_one_authored_page_becomes_one_page_per_frame(self):
        self.assertEqual(
            _page_ids(_lecture()),
            ["before", "anim-1", "anim-2", "anim-3", "after"],
        )

    def test_every_frame_repeats_the_pages_text(self):
        frames = _pages(_lecture())[1:4]
        self.assertEqual({page.title for page in frames}, {"Commit logging 先写完整新状态"})
        slides = {
            block.content
            for page in frames
            for block in page.blocks
            if block.kind == "slide"
        }
        self.assertEqual(len(slides), 1)

    def test_only_the_image_source_changes(self):
        frames = _pages(_lecture())[1:4]
        self.assertEqual(
            [_images(page)[0].content["src"] for page in frames], list(SOURCES)
        )

    def test_the_figure_lands_where_frames_was_called(self):
        frame = _pages(_lecture())[1]
        self.assertEqual([block.kind for block in frame.blocks],
                         ["slide", "image", "notes"])

    def test_shared_options_ride_every_frame(self):
        for page in _pages(_lecture())[1:4]:
            image = _images(page)[0]
            self.assertEqual(image.content["width"], "900px")
            self.assertEqual(image.content["caption"], "先写 log，再改数据")
            self.assertEqual(image.footnotes, ("Saltzer & Kaashoek §9.3",))

    def test_bubbles_ride_every_frame(self):
        def body(p):
            p.title("T")
            p.frames("a.png", "b.png").annotate("look here", at="center")

        lec = Lecture(id="lec-a", title="A")
        lec.page(id="anim", body=body)
        for page in _pages(lec.build()):
            self.assertEqual(_images(page)[0].annotations[0].text, "look here")

    def test_frame_group_records_the_position(self):
        groups = [page.frame_group for page in _pages(_lecture())[1:4]]
        self.assertEqual(
            [(g.id, g.index, g.total) for g in groups],
            [("anim", 1, 3), ("anim", 2, 3), ("anim", 3, 3)],
        )

    def test_a_page_without_frames_is_untouched(self):
        page = _pages(_lecture())[0]
        self.assertEqual(page.id, "before")
        self.assertIsNone(page.frame_group)

    def test_a_single_frame_is_still_a_group(self):
        def body(p):
            p.title("T")
            p.frames("only.png")

        lec = Lecture(id="lec-1", title="One")
        lec.page(id="anim", body=body)
        [page] = _pages(lec.build())
        self.assertEqual(page.id, "anim-1")
        self.assertEqual(page.frame_group.total, 1)

    def test_frame_groups_maps_a_group_to_its_deck_positions(self):
        self.assertEqual(
            model.frame_groups(_pages(_lecture())), {"anim": [1, 2, 3]}
        )


class FramesValidationTest(unittest.TestCase):
    def _build(self, body):
        lec = Lecture(id="lec-e", title="E")
        lec.page(id="anim", body=body)
        return lec.build()

    def test_an_animation_needs_at_least_one_image(self):
        def body(p):
            p.title("T")
            p.frames()

        with self.assertRaises(ValidationError) as caught:
            self._build(body)
        self.assertIn("at least one image", str(caught.exception))

    def test_a_blank_source_is_rejected(self):
        def body(p):
            p.title("T")
            p.frames("a.png", "  ")

        with self.assertRaises(ValidationError):
            self._build(body)

    def test_a_page_carries_at_most_one_animation(self):
        def body(p):
            p.title("T")
            p.frames("a.png", "b.png")
            p.frames("c.png", "d.png")

        with self.assertRaises(ValidationError) as caught:
            self._build(body)
        self.assertIn("at most one frames block", str(caught.exception))

    def test_the_expanded_ast_carries_no_frames_block(self):
        # Every renderer's block table is keyed on model.BLOCK_KINDS, which the
        # transient frames kind is deliberately absent from.
        self.assertNotIn("frames", model.BLOCK_KINDS)
        kinds = {block.kind for page in _pages(_lecture()) for block in page.blocks}
        self.assertNotIn("frames", kinds)


def _book_lecture(**page_kwargs):
    lec = Lecture(id="lec-b", title="Lecture B")

    def body(p):
        p.title("Commit logging")
        p.prose("正文：先写 log，再改数据。")
        p.frames(*SOURCES, caption="commit logging 的三个时刻", **page_kwargs)

    lec.page(id="anim", body=body)
    return lec.build()


def _render_book(lecture):
    tmp = Path(tempfile.mkdtemp())
    for name in SOURCES:
        (tmp / name).write_bytes(b"png")  # AssetCopier insists images exist
    book = BookModel(
        title="T", author=None, subtitle=None, preface=None,
        lectures=(lecture,), asset_roots={lecture.id: tmp},
    )
    stderr = io.StringIO()
    with contextlib.redirect_stderr(stderr):
        LatexRenderer().render(book, tmp)
    return (tmp / "chapters" / f"{lecture.id}.tex").read_text(), stderr.getvalue()


class BookTest(unittest.TestCase):
    def test_the_book_prints_only_the_last_frame(self):
        tex, _ = _render_book(_book_lecture())
        self.assertIn("log-3", tex)
        self.assertNotIn("log-1", tex)
        self.assertNotIn("log-2", tex)

    def test_the_animation_is_one_section_with_one_caption(self):
        tex, _ = _render_book(_book_lecture())
        self.assertEqual(tex.count(r"\section{Commit logging}"), 1)
        self.assertEqual(tex.count("commit logging 的三个时刻"), 1)

    def test_earlier_frames_earn_no_todo_box(self):
        tex, _ = _render_book(_book_lecture())
        self.assertNotIn(r"\booktodo", tex)

    def test_an_animation_counts_as_one_unit_of_coverage(self):
        book = BookModel(
            title="T", author=None, subtitle=None, preface=None,
            lectures=(_book_lecture(),), asset_roots={"lec-b": Path(".")},
        )
        self.assertEqual(coverage(book), [("lec-b", 1, 1)])

    def test_a_ref_names_the_printed_figure_once(self):
        tex, stderr = _render_book(_book_lecture(ref="commit-log"))
        self.assertEqual(tex.count(r"\label{fig:lec-b-commit-log}"), 1)
        self.assertEqual(stderr, "")

    def test_the_deck_keeps_every_frame(self):
        self.assertEqual(_page_ids(_book_lecture()), ["anim-1", "anim-2", "anim-3"])


class PageSelectionTest(unittest.TestCase):
    def test_the_authored_page_id_selects_the_whole_animation(self):
        out = model.select_pages(_lecture(), "anim")
        self.assertEqual(_page_ids(out), ["anim-1", "anim-2", "anim-3"])

    def test_a_frame_id_selects_that_frame(self):
        out = model.select_pages(_lecture(), "anim-2")
        self.assertEqual(_page_ids(out), ["anim-2"])

    def test_a_number_is_still_a_deck_index(self):
        out = model.select_pages(_lecture(), "3")
        self.assertEqual(_page_ids(out), ["anim-2"])

    def test_an_unknown_id_still_fails(self):
        with self.assertRaises(ValidationError):
            model.select_pages(_lecture(), "anim-9")


def _tree_pages(node_list):
    out = []
    for node in node_list:
        if node["type"] == "page":
            out.append(node)
        else:
            out.extend(_tree_pages(node["children"]))
    return out


class OutlineTest(unittest.TestCase):
    def test_an_animation_is_one_outline_row(self):
        rows = _tree_pages(build_data(_lecture())["tree"])
        self.assertEqual([row["id"] for row in rows], ["before", "anim-1", "after"])

    def test_the_row_carries_the_number_of_slides_it_stands_for(self):
        rows = {row["id"]: row for row in _tree_pages(build_data(_lecture())["tree"])}
        self.assertEqual(rows["anim-1"]["frames"], 3)
        self.assertNotIn("frames", rows["before"])

    def test_every_frame_is_still_a_slide(self):
        pages = build_data(_lecture())["pages"]
        self.assertEqual(
            [page["id"] for page in pages],
            ["before", "anim-1", "anim-2", "anim-3", "after"],
        )

    def test_the_printed_outline_labels_the_row_with_one_number(self):
        html = build_outline_html(_lecture())
        self.assertIn('<span class="page-number">2</span>', html)

    def test_an_animation_advances_the_shown_number_by_one(self):
        # before / anim (3 frames) / after are shown as 1 / 2 / 3 — an animation
        # is one slide to a reader, so the page after it is 3, not 5.
        html = build_outline_html(_lecture())
        self.assertIn('<span class="page-number">3</span>', html)
        self.assertEqual(html.count("page-number"), 3)

    def test_the_row_still_links_to_the_physical_slide(self):
        # The label compresses; the link must not — it addresses a real slide,
        # so "after" stays deck slide 5 (0-based 4) even though it shows as 3.
        html = build_outline_html(_lecture())
        self.assertIn(f'href="{SLIDE_LINK_PREFIX}1"', html)   # anim-1
        self.assertIn(f'href="{SLIDE_LINK_PREFIX}4"', html)   # after

    def test_a_pruned_animation_spans_only_the_frames_that_survived(self):
        pruned = model.select_pages(_lecture(), "anim-2,anim-3")
        rows = {row["id"]: row for row in _tree_pages(build_data(pruned)["tree"])}
        self.assertEqual(rows["anim-2"]["frames"], 2)


class ShownNumberTest(unittest.TestCase):
    """An animation is one slide to a reader: every frame shows one number."""

    def _slides(self):
        return build_marp_markdown(_lecture()).split("\n\n---\n\n")

    def test_the_first_frame_carries_the_number(self):
        slides = self._slides()
        self.assertNotIn("_paginate: hold", slides[1])

    def test_later_frames_hold_it(self):
        slides = self._slides()
        self.assertIn("<!-- _paginate: hold -->", slides[2])
        self.assertIn("<!-- _paginate: hold -->", slides[3])

    def test_the_page_after_an_animation_advances_again(self):
        self.assertNotIn("_paginate: hold", self._slides()[4])

    def test_a_deck_without_animations_never_holds(self):
        lec = Lecture(id="plain", title="Plain")
        lec.page(id="one", body=lambda p: (p.title("One"), p.slide("a")))
        lec.page(id="two", body=lambda p: (p.title("Two"), p.slide("b")))
        self.assertNotIn("_paginate: hold", build_marp_markdown(lec.build()))

    def test_the_json_carries_the_shown_number_per_page(self):
        pages = build_data(_lecture())["pages"]
        self.assertEqual([page["number"] for page in pages], [1, 2, 2, 2, 3])


class SlideNumbersTest(unittest.TestCase):
    def test_pages_without_animations_number_themselves(self):
        pages = _pages(_lecture())
        self.assertEqual(model.slide_numbers(pages), [1, 2, 2, 2, 3])

    def test_an_empty_deck_is_empty(self):
        self.assertEqual(model.slide_numbers([]), [])

    def test_a_pruned_animation_still_counts_once(self):
        pruned = _pages(model.select_pages(_lecture(), "anim-2,anim-3,after"))
        self.assertEqual(model.slide_numbers(pruned), [1, 1, 2])


class RevealTest(unittest.TestCase):
    def _slides(self, **kwargs):
        return build_marp_markdown(_lecture(), **kwargs).split("\n\n---\n\n")

    def test_the_first_frame_steps_through_its_blocks(self):
        slides = self._slides(reveal=True)
        self.assertIn("data-reveal", slides[1])

    def test_later_frames_arrive_fully_lit(self):
        slides = self._slides(reveal=True)
        self.assertNotIn("data-reveal", slides[2])
        self.assertNotIn("data-reveal", slides[3])

    def test_ordinary_pages_are_unaffected(self):
        slides = self._slides(reveal=True)
        self.assertIn("data-reveal", slides[4])

    def test_without_reveal_nothing_steps(self):
        self.assertNotIn("data-reveal", "".join(self._slides()))

    def test_the_speaker_note_rides_every_frame(self):
        slides = self._slides()
        self.assertEqual(sum(NOTE in slide for slide in slides[1:4]), 3)


def _cited_lecture():
    lec = Lecture(id="lec-c", title="Lecture C")

    def before(p):
        p.title("Before")
        p.slide("first")

    def anim(p):
        p.title("Animated")
        p.frames("a.png", "b.png", "c.png")
        p.cite(title="Principles of Computer System Design", year="2009")

    def after(p):
        p.title("After")
        p.slide("last")
        p.cite(title="Principles of Computer System Design", year="2009")

    lec.page(id="before", body=before)
    lec.page(id="anim", body=anim)
    lec.page(id="after", body=after)
    return lec.build()


class CitationBackrefTest(unittest.TestCase):
    def test_an_animation_reports_one_slide_not_every_frame(self):
        [entry] = references.collect_citations(_pages(_cited_lecture()))
        # Cited on the animation and again on the page after it: the animation
        # counts once, and shows as one number, so the pages after it shift up.
        self.assertEqual(entry.pages, (2, 3))

    def test_the_deck_reference_page_shows_the_folded_backref(self):
        pages = build_data(_cited_lecture())["pages"]
        text = pages[-1]["blocks"][0]["content"]
        self.assertIn("(P2, P3)", text)


class SerializeTest(unittest.TestCase):
    def test_the_ast_records_each_frames_position(self):
        pages = [
            node
            for child in lecture_to_dict(_lecture())["children"]
            for node in child["children"]
        ]
        by_id = {page["id"]: page for page in pages}
        self.assertEqual(
            by_id["anim-2"]["frame_group"], {"id": "anim", "index": 2, "total": 3}
        )
        self.assertIsNone(by_id["before"]["frame_group"])


def _tail_lecture(*, sources=SOURCES, gap=False):
    """before / an animation with a trailing verdict / nothing after."""
    lec = Lecture(id="lec-t", title="Lecture T")

    def before(p):
        p.title("Before")
        p.slide("first")

    def body(p):
        p.title("Case")
        if gap:
            p.gap()
        p.slide("setup")                                   # every frame
        p.frames(*sources)
        p.highlight("verdict").footnote("来源：某论文")     # last frame only

    lec.page(id="before", body=before)
    lec.page(id="anim", body=body)
    return lec.build()


class AfterFramesTest(unittest.TestCase):
    """Blocks written after ``p.frames(...)`` belong to the finished picture."""

    def _slides(self, lecture, **kwargs):
        return build_marp_markdown(lecture, **kwargs).split("\n\n---\n\n")

    def test_tail_blocks_are_marked_on_every_frame(self):
        for page in _pages(_tail_lecture())[1:4]:
            marks = {block.kind: block.after_frames for block in page.blocks}
            self.assertEqual(
                marks, {"slide": False, "image": False, "highlight": True}
            )

    def test_earlier_frames_hold_the_tail_block(self):
        slides = self._slides(_tail_lecture())
        for frame in (slides[1], slides[2]):
            self.assertIn("lk-held", frame)
            self.assertIn("verdict", frame)     # rendered — invisible, in place
        self.assertNotIn("lk-held", slides[3])
        self.assertIn("verdict", slides[3])

    def test_a_held_footnote_waits_for_the_last_frame(self):
        slides = self._slides(_tail_lecture())
        self.assertNotIn("来源", slides[1])
        self.assertIn("来源", slides[3])

    def test_pre_animation_blocks_still_ride_every_frame(self):
        slides = self._slides(_tail_lecture())
        for frame in slides[1:4]:
            self.assertIn("setup", frame)
            self.assertNotIn('class="lk-held">\n\nsetup', frame)

    def test_reveal_steps_each_half_where_it_first_appears(self):
        slides = self._slides(_tail_lecture(), reveal=True)
        self.assertEqual(slides[1].count("data-reveal"), 2)   # setup + figure
        self.assertIn("lk-held", slides[1])                   # verdict hidden
        self.assertNotIn("data-reveal", slides[2])            # middle: lit
        self.assertEqual(slides[3].count("data-reveal"), 1)   # verdict steps
        self.assertNotIn("lk-held", slides[3])

    def test_a_single_frame_animation_steps_everything_at_once(self):
        slides = self._slides(_tail_lecture(sources=("only.png",)), reveal=True)
        self.assertNotIn("lk-held", slides[1])
        # setup + figure + verdict: the single frame is first and last at once.
        self.assertEqual(slides[1].count("data-reveal"), 3)

    def test_a_held_block_keeps_its_seam_in_the_gap_flow(self):
        slides = self._slides(_tail_lecture(gap=True))
        self.assertIn('class="lk-gap-block lk-held"', slides[1])
        self.assertEqual(
            slides[1].count("lk-gap-spacer"), slides[3].count("lk-gap-spacer")
        )

    def test_the_book_prints_the_tail_with_the_finished_picture(self):
        lec = Lecture(id="lec-b", title="Lecture B")

        def body(p):
            p.title("Case")
            p.prose("正文。")
            p.frames(*SOURCES, caption="三个时刻")
            p.highlight("verdict")

        lec.page(id="anim", body=body)
        tex, _ = _render_book(lec.build())
        self.assertEqual(tex.count("verdict"), 1)

    def test_pptx_drops_held_blocks_from_earlier_frames(self):
        import tempfile as _tempfile

        from pptx import Presentation
        from lecturekit.renderers.pptx import PptxRenderer

        tmp = Path(_tempfile.mkdtemp())
        png = (
            b"\x89PNG\r\n\x1a\n" + bytes.fromhex(
                "0000000d49484452000000010000000108060000001f15c489"
                "0000000d4944415478da63fccff0bf1e00057c02b5b0f75db6"
                "0000000049454e44ae426082"
            )
        )
        for name in SOURCES:
            (tmp / name).write_bytes(png)
        out = PptxRenderer(asset_root=tmp).render(_tail_lecture(), tmp)
        prs = Presentation(str(out))
        texts = [
            "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            for slide in prs.slides
        ]
        self.assertNotIn("verdict", texts[1])   # frame 1
        self.assertNotIn("verdict", texts[2])   # frame 2
        self.assertIn("verdict", texts[3])      # last frame

    def test_a_side_image_is_exempt_from_holding(self):
        lec = Lecture(id="lec-s", title="S")

        def body(p):
            p.title("T")
            p.frames("a.png", "b.png")
            p.side_image("bg.png")

        lec.page(id="anim", body=body)
        pages = _pages(lec.build())
        side = next(b for b in pages[0].blocks if b.kind == "side_image")
        self.assertTrue(side.after_frames)
        self.assertFalse(model.block_held(pages[0], side))

    def test_the_ast_records_the_mark(self):
        pages = [
            node for node in lecture_to_dict(_tail_lecture())["children"]
            if node["type"] == "page"
        ]
        by_id = {page["id"]: page for page in pages}
        kinds = {
            block["kind"]: block["after_frames"]
            for block in by_id["anim-2"]["blocks"]
        }
        self.assertEqual(
            kinds, {"slide": False, "image": False, "highlight": True}
        )


class InspectTest(unittest.TestCase):
    def test_inspect_lists_an_animation_once_under_its_authored_id(self):
        tree = format_tree(_lecture())
        self.assertIn("(anim) [3 frames]", tree)
        self.assertNotIn("anim-2", tree)


if __name__ == "__main__":
    unittest.main()
