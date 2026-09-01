"""PPTX embeds SVG figures by rasterizing them first."""

import contextlib
import io
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from lecturekit import Lecture, rasterize
from lecturekit.renderers.pptx import PptxRenderer

# A minimal SVG every backend can render: one blue square, explicitly sized.
_SVG = """<svg xmlns="http://www.w3.org/2000/svg" width="120" height="60"
     viewBox="0 0 120 60"><rect width="120" height="60" fill="#1769C2"/></svg>
"""


def pictures(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


class SvgRasterizerTest(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())
        self.svg = self.tmp / "square.svg"
        self.svg.write_text(_SVG)

    def test_converts_to_a_png(self):
        if rasterize.find_backend() is None:
            self.skipTest("no SVG backend installed")
        png = rasterize.SvgRasterizer(self.tmp / "out").png(self.svg)
        self.assertIsNotNone(png)
        self.assertEqual(png.suffix, ".png")
        # PNG signature, so this is a real raster and not a copied SVG.
        self.assertEqual(png.read_bytes()[:8], b"\x89PNG\r\n\x1a\n")

    def test_converts_each_source_once(self):
        calls = []

        def convert(src, dest, scale):
            calls.append(src)
            dest.write_bytes(b"png")

        r = rasterize.SvgRasterizer(self.tmp / "out")
        r._backend = ("fake", convert)
        first, second = r.png(self.svg), r.png(self.svg)
        self.assertEqual(first, second)
        self.assertEqual(calls, [self.svg])

    def test_same_basename_from_two_directories_does_not_collide(self):
        other = self.tmp / "sub"
        other.mkdir()
        twin = other / "square.svg"
        twin.write_text(_SVG)

        r = rasterize.SvgRasterizer(self.tmp / "out")
        r._backend = ("fake", lambda src, dest, scale: dest.write_bytes(b"png"))
        self.assertNotEqual(r.png(self.svg), r.png(twin))

    def test_no_backend_reports_no_png(self):
        r = rasterize.SvgRasterizer(self.tmp / "out")
        r._backend = None
        self.assertIsNone(r.png(self.svg))
        self.assertIsNone(r.backend)

    def test_a_backend_that_fails_costs_only_that_figure(self):
        def explode(src, dest, scale):
            raise RuntimeError("bad svg")

        r = rasterize.SvgRasterizer(self.tmp / "out")
        r._backend = ("fake", explode)
        self.assertIsNone(r.png(self.svg))


class PptxSvgTest(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())
        (self.tmp / "assets").mkdir()
        (self.tmp / "assets" / "square.svg").write_text(_SVG)

    def deck(self):
        lecture = Lecture(id="lec", title="Deck")

        def body(p):
            p.title("A vector figure")
            p.image("assets/square.svg", caption="正方形", width_px=240)

        lecture.page(id="fig", body=body)
        out = PptxRenderer(asset_root=self.tmp).render(
            lecture.build(), self.tmp / "out"
        )
        return Presentation(str(out))

    def test_svg_figure_is_embedded(self):
        if rasterize.find_backend() is None:
            self.skipTest("no SVG backend installed")
        slide = self.deck().slides[0]
        self.assertEqual(len(pictures(slide)), 1)

    def test_aspect_ratio_survives_the_conversion(self):
        if rasterize.find_backend() is None:
            self.skipTest("no SVG backend installed")
        pic = pictures(self.deck().slides[0])[0]
        # The source SVG is 120x60, so the picture must stay 2:1.
        self.assertAlmostEqual(pic.width / pic.height, 2.0, places=2)

    def test_a_figure_that_cannot_be_converted_keeps_its_caption_and_warns(self):
        if rasterize.find_backend() is None:
            self.skipTest("no SVG backend installed")
        (self.tmp / "assets" / "broken.svg").write_text("not an svg at all")
        lecture = Lecture(id="lec", title="Deck")

        def body(p):
            p.title("A broken figure")
            p.image("assets/broken.svg", caption="仍然有 caption")

        lecture.page(id="fig", body=body)
        stderr = io.StringIO()
        with contextlib.redirect_stderr(stderr):
            out = PptxRenderer(asset_root=self.tmp).render(
                lecture.build(), self.tmp / "out"
            )
        slide = Presentation(str(out)).slides[0]
        self.assertEqual(pictures(slide), [])
        text = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        self.assertIn("仍然有 caption", text)
        self.assertIn("broken.svg", stderr.getvalue())


if __name__ == "__main__":
    unittest.main()
